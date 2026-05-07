from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Inches, Pt

from deckbridge.renderers.common.style_resolver import (
    resolve_chart_theme,
    resolve_series_color,
    resolve_series_dash,
    resolve_series_width,
)
from deckbridge.renderers.gslides.utils import GSLIDES_LINE_DASH_MAP, hex_to_slides_rgb, inches_to_emu
from deckbridge.renderers.pptx.utils import PPTX_DASH_MAP, hex_to_rgb255


def render_dash_legend(ctx, slot_key, slot, slide):
    source_key = slot["source"]

    block = slide["content"].get(source_key)

    if not block:
        return

    chart = block.chart

    chart_theme = resolve_chart_theme(
        ctx.theme,
        ctx.layout_spec.name,
    )

    if ctx.backend == "pptx":
        _render_dash_legend_pptx(
            ctx,
            slot,
            chart,
            chart_theme,
        )

    elif ctx.backend == "gslides":
        _render_dash_legend_gslides(
            ctx,
            slot_key,
            slot,
            chart,
            chart_theme,
        )


def _render_dash_legend_pptx(ctx, slot, chart, chart_theme):
    x = slot["x"]
    y = slot["y"]

    line_w = 0.35
    row_h = 0.25

    for i, series in enumerate(chart.series):
        color = resolve_series_color(series, i, chart_theme)
        dash = resolve_series_dash(series, chart_theme)
        width = resolve_series_width(series, chart_theme)

        y_i = y + i * row_h

        line = ctx.slide_obj.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(x),
            Inches(y_i),
            Inches(x + line_w),
            Inches(y_i),
        )

        line.line.color.rgb = hex_to_rgb255(color)
        line.line.width = Pt(width)
        line.shadow.inherit = False

        if dash != "solid":
            line.line.dash_style = PPTX_DASH_MAP[dash]

        textbox = ctx.slide_obj.shapes.add_textbox(
            Inches(x + 0.45),
            Inches(y_i - 0.15),
            Inches(1.5),
            Inches(0.3),
        )

        textbox.text_frame.text = series.get("name", series["column"])
        textbox.text_frame.paragraphs[0].font.size = Pt(12)


def _render_dash_legend_gslides(ctx, slot_key, slot, chart, chart_theme):
    requests = []

    x = slot["x"]
    y = slot["y"]

    line_w = 0.35
    row_h = 0.25

    for i, series in enumerate(chart.series):
        color = resolve_series_color(series, i, chart_theme)
        dash = resolve_series_dash(series, chart_theme)
        width = resolve_series_width(series, chart_theme)

        y_i = y + i * row_h

        line_id = f"{slot_key}_line_{i}_{ctx.page_id}"

        requests.append(
            {
                "createLine": {
                    "objectId": line_id,
                    "category": "STRAIGHT",
                    "elementProperties": {
                        "pageObjectId": ctx.page_id,
                        "size": {
                            "width": {
                                "magnitude": inches_to_emu(line_w),
                                "unit": "EMU",
                            },
                            "height": {
                                "magnitude": 0,
                                "unit": "EMU",
                            },
                        },
                        "transform": {
                            "scaleX": 1,
                            "scaleY": 1,
                            "translateX": inches_to_emu(x),
                            "translateY": inches_to_emu(y_i),
                            "unit": "EMU",
                        },
                    },
                }
            }
        )

        requests.append(
            {
                "updateLineProperties": {
                    "objectId": line_id,
                    "lineProperties": {
                        "weight": {
                            "magnitude": width,
                            "unit": "PT",
                        },
                        "dashStyle": GSLIDES_LINE_DASH_MAP[dash],
                        "lineFill": {"solidFill": {"color": {"rgbColor": hex_to_slides_rgb(color)}}},
                    },
                    "fields": "weight,dashStyle,lineFill",
                }
            }
        )

        text_id = f"{slot_key}_text_{i}_{ctx.page_id}"

        requests.append(
            {
                "createShape": {
                    "objectId": text_id,
                    "shapeType": "TEXT_BOX",
                    "elementProperties": {
                        "pageObjectId": ctx.page_id,
                        "size": {
                            "height": {
                                "magnitude": inches_to_emu(0.3),
                                "unit": "EMU",
                            },
                            "width": {
                                "magnitude": inches_to_emu(1.5),
                                "unit": "EMU",
                            },
                        },
                        "transform": {
                            "scaleX": 1,
                            "scaleY": 1,
                            "translateX": inches_to_emu(x + 0.45),
                            "translateY": inches_to_emu(y_i - 0.15),
                            "unit": "EMU",
                        },
                    },
                }
            }
        )

        requests.append(
            {
                "insertText": {
                    "objectId": text_id,
                    "text": series.get("name", series["column"]),
                }
            }
        )

        requests.append(
            {
                "updateTextStyle": {
                    "objectId": text_id,
                    "textRange": {"type": "ALL"},
                    "style": {
                        "fontSize": {
                            "magnitude": 12,
                            "unit": "PT",
                        },
                        "bold": False,
                    },
                    "fields": "fontSize,bold",
                }
            }
        )

        requests.append(
            {
                "updateParagraphStyle": {
                    "objectId": text_id,
                    "textRange": {"type": "ALL"},
                    "style": {
                        "alignment": "START",
                    },
                    "fields": "alignment",
                }
            }
        )

        requests.append(
            {
                "updateShapeProperties": {
                    "objectId": text_id,
                    "shapeProperties": {"contentAlignment": "MIDDLE"},
                    "fields": "contentAlignment",
                }
            }
        )

    ctx.slides_service.presentations().batchUpdate(
        presentationId=ctx.presentation_id,
        body={"requests": requests},
    ).execute()
